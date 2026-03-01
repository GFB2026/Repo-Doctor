#!/usr/bin/env python3
"""
doc-intel.py - Document Intelligence Layer for Project Doctor
Extracts, OCRs, indexes, and searches content from ANY file type across all projects.

Usage:
    python doc-intel.py scan <project-name>              # Extract & index everything
    python doc-intel.py scan <project-name> --deep       # Include archives + database rows
    python doc-intel.py search <query>                    # Full-text search ALL projects
    python doc-intel.py search <query> --project <n>   # Search within one project
    python doc-intel.py info <project-name>               # Document stats for a project
    python doc-intel.py entities <project-name> [type]    # Extracted entities (dates, $, etc)
    python doc-intel.py summary                           # Cross-project document overview
    python doc-intel.py deps-check                        # Check which extractors are installed
    python doc-intel.py rebuild                            # Rebuild full-text search index

Core extractors (pip install):
    pip install pdfplumber python-docx openpyxl python-pptx --break-system-packages

Extended extractors (pip install as needed):
    pip install xlrd olefile striprtf python-calamine --break-system-packages
    pip install extract-msg eml-parser --break-system-packages
    pip install Pillow pytesseract --break-system-packages
    pip install chardet --break-system-packages

System tools (optional):
    tesseract-ocr  - OCR for scanned PDFs/images
    poppler        - PDF to image conversion (pdftoppm)
    libreoffice    - Fallback for .doc, .odt, .ods, .odp
"""

import sqlite3
import json
import os
import sys
import re
import csv
import hashlib
import struct
import zipfile
import tarfile
import io
from datetime import datetime, timezone
from pathlib import Path


# ─── Configuration ────────────────────────────────────────────────────────────

DB_DIR = Path.home() / ".repo-doctor"
DB_PATH = DB_DIR / "registry.db"

MAX_TEXT_SIZE = 10 * 1024 * 1024   # 10MB text limit per file
MAX_FILE_SIZE = 500 * 1024 * 1024  # 500MB skip threshold
ARCHIVE_MAX_FILES = 200            # Max files to extract from a single archive
DB_MAX_ROWS = 10000                # Max rows to extract from a database table

# Every file type we can handle, mapped to its extraction strategy
SUPPORTED_TYPES = {
    # ── Documents ──
    ".pdf":   "pdf",
    ".docx":  "docx",
    ".doc":   "doc_legacy",
    ".rtf":   "rtf",
    ".odt":   "odt",
    ".pages": "pages",

    # ── Spreadsheets ──
    ".xlsx":  "xlsx",
    ".xls":   "xls",
    ".xlsm":  "xlsx",
    ".xlsb":  "xlsb",
    ".ods":   "ods",
    ".numbers": "numbers",

    # ── Presentations ──
    ".pptx":  "pptx",
    ".ppt":   "ppt_legacy",
    ".odp":   "odp",
    ".key":   "keynote",

    # ── Databases ──
    ".db":    "sqlite",
    ".sqlite": "sqlite",
    ".sqlite3": "sqlite",
    ".mdb":   "access",
    ".accdb": "access",

    # ── Email ──
    ".eml":   "eml",
    ".msg":   "msg",
    ".mbox":  "mbox",
    ".pst":   "pst",

    # ── Plain Text / Code / Config ──
    ".txt":   "text",
    ".md":    "text",
    ".markdown": "text",
    ".rst":   "text",
    ".csv":   "csv",
    ".tsv":   "tsv",
    ".json":  "json_file",
    ".jsonl": "jsonl",
    ".ndjson": "jsonl",
    ".xml":   "xml",
    ".yaml":  "text",
    ".yml":   "text",
    ".toml":  "text",
    ".ini":   "text",
    ".cfg":   "text",
    ".conf":  "text",
    ".env":   "text",
    ".log":   "text",
    ".sql":   "text",
    ".html":  "html",
    ".htm":   "html",
    ".tex":   "text",
    ".bib":   "text",

    # ── Code (index for search, not audit) ──
    ".py":    "text",
    ".js":    "text",
    ".ts":    "text",
    ".jsx":   "text",
    ".tsx":   "text",
    ".java":  "text",
    ".c":     "text",
    ".cpp":   "text",
    ".cs":    "text",
    ".go":    "text",
    ".rs":    "text",
    ".rb":    "text",
    ".php":   "text",
    ".swift": "text",
    ".kt":    "text",
    ".r":     "text",
    ".R":     "text",
    ".sh":    "text",
    ".bash":  "text",
    ".bat":   "text",
    ".ps1":   "text",
    ".psm1":  "text",
    ".css":   "text",
    ".scss":  "text",
    ".less":  "text",

    # ── Data / Notebooks ──
    ".ipynb": "jupyter",
    ".parquet": "parquet",

    # ── Images (OCR) ──
    ".png":   "image",
    ".jpg":   "image",
    ".jpeg":  "image",
    ".tif":   "image",
    ".tiff":  "image",
    ".bmp":   "image",
    ".webp":  "image",
    ".heic":  "image",

    # ── Archives (recurse into) ──
    ".zip":   "archive_zip",
    ".tar":   "archive_tar",
    ".gz":    "archive_gz",
    ".tgz":   "archive_tar",
    ".bz2":   "archive_tar",
    ".7z":    "archive_7z",
    ".rar":   "archive_rar",

    # ── Calendar / Contacts ──
    ".ics":   "ical",
    ".vcf":   "vcard",
}

CATEGORY_MAP = {
    "pdf": "document", "docx": "document", "doc_legacy": "document", "rtf": "document",
    "odt": "document", "pages": "document",
    "xlsx": "spreadsheet", "xls": "spreadsheet", "xlsb": "spreadsheet", "ods": "spreadsheet",
    "numbers": "spreadsheet", "csv": "spreadsheet", "tsv": "spreadsheet",
    "pptx": "presentation", "ppt_legacy": "presentation", "odp": "presentation",
    "keynote": "presentation",
    "sqlite": "database", "access": "database",
    "eml": "email", "msg": "email", "mbox": "email", "pst": "email",
    "text": "text", "html": "text", "xml": "text", "json_file": "data",
    "jsonl": "data", "jupyter": "notebook", "parquet": "data",
    "image": "image",
    "archive_zip": "archive", "archive_tar": "archive", "archive_gz": "archive",
    "archive_7z": "archive", "archive_rar": "archive",
    "ical": "calendar", "vcard": "contact",
}


# ─── Database Schema ─────────────────────────────────────────────────────────

DOC_SCHEMA = """
CREATE TABLE IF NOT EXISTS documents (
    id              INTEGER PRIMARY KEY AUTOINCREMENT,
    repo_id         INTEGER NOT NULL,
    file_path       TEXT NOT NULL,
    file_name       TEXT NOT NULL,
    file_type       TEXT,
    file_category   TEXT,
    file_hash       TEXT,
    file_size       INTEGER,
    page_count      INTEGER,
    word_count      INTEGER,
    char_count      INTEGER,
    extracted_text   TEXT,
    extraction_method TEXT,
    is_searchable   INTEGER DEFAULT 0,
    is_ocrd         INTEGER DEFAULT 0,
    language        TEXT,
    modified_date   TEXT,
    scanned_at      TEXT NOT NULL,
    metadata        TEXT,
    UNIQUE(repo_id, file_path)
);

CREATE VIRTUAL TABLE IF NOT EXISTS doc_search USING fts5(
    file_path,
    file_name,
    extracted_text,
    repo_name,
    repo_id UNINDEXED,
    doc_id UNINDEXED
);

CREATE TABLE IF NOT EXISTS doc_entities (
    id              INTEGER PRIMARY KEY AUTOINCREMENT,
    doc_id          INTEGER NOT NULL REFERENCES documents(id) ON DELETE CASCADE,
    repo_id         INTEGER NOT NULL,
    entity_type     TEXT NOT NULL,
    entity_value    TEXT NOT NULL,
    context         TEXT,
    confidence      REAL DEFAULT 1.0
);

CREATE INDEX IF NOT EXISTS idx_doc_repo ON documents(repo_id);
CREATE INDEX IF NOT EXISTS idx_doc_hash ON documents(file_hash);
CREATE INDEX IF NOT EXISTS idx_doc_cat ON documents(file_category);
CREATE INDEX IF NOT EXISTS idx_entity_type ON doc_entities(entity_type);
CREATE INDEX IF NOT EXISTS idx_entity_repo ON doc_entities(repo_id);
"""


# ─── Helpers ──────────────────────────────────────────────────────────────────

def get_db():
    DB_DIR.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")
    return conn

from contextlib import contextmanager

@contextmanager
def db_session(schema=None):
    """Context manager for DB connections: connect, init schema, commit/rollback, close."""
    conn = get_db()
    if schema:
        conn.executescript(schema)
    try:
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()

def now_iso():
    return datetime.now(timezone.utc).isoformat()

def file_hash(filepath):
    h = hashlib.md5()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()

# Pre-compiled regex for word counting (Phase 5C: O(1) memory)
_WORD_RE = re.compile(r'\S+')

def count_words(text):
    if not text:
        return 0
    return sum(1 for _ in _WORD_RE.finditer(text))

def truncate(text, max_len=MAX_TEXT_SIZE):
    if text and len(text) > max_len:
        return text[:max_len] + f"\n[TRUNCATED at {max_len} chars]"
    return text

def safe_read_text(filepath, encoding=None):
    """Read text file with smarter encoding detection."""
    if encoding:
        try:
            return Path(filepath).read_text(encoding=encoding, errors="replace")
        except Exception:
            pass
    # Try UTF-8 strict first (most common case)
    try:
        return Path(filepath).read_text(encoding="utf-8")
    except (UnicodeDecodeError, ValueError):
        pass
    # Chardet on first 10KB, then fallback to latin-1
    try:
        import chardet
        raw = Path(filepath).read_bytes()
        detected = chardet.detect(raw[:10240])
        enc = detected.get("encoding") or "latin-1"
        return raw.decode(enc, errors="replace")
    except ImportError:
        pass
    return Path(filepath).read_bytes().decode("latin-1", errors="replace")


# Pre-compiled HTML stripping patterns (Phase 4C)
_RE_SCRIPT = re.compile(r"<script[^>]*>.*?</script>", re.DOTALL | re.IGNORECASE)
_RE_STYLE = re.compile(r"<style[^>]*>.*?</style>", re.DOTALL | re.IGNORECASE)
_RE_TAGS = re.compile(r"<[^>]+>")
_RE_WHITESPACE = re.compile(r"\s+")


def strip_html(text):
    """Strip HTML tags from text using pre-compiled patterns."""
    text = _RE_SCRIPT.sub("", text)
    text = _RE_STYLE.sub("", text)
    text = _RE_TAGS.sub(" ", text)
    return _RE_WHITESPACE.sub(" ", text).strip()


def format_email_headers(from_, to_, date_, subject_):
    """Format common email headers."""
    return [f"From: {from_}", f"To: {to_}", f"Date: {date_}", f"Subject: {subject_}", ""]


# ─── Extractors ───────────────────────────────────────────────────────────────

def extract_pdf(fp):
    try:
        import pdfplumber
        with pdfplumber.open(fp) as pdf:
            pages = len(pdf.pages)
            text = "\n".join(p.extract_text() or "" for p in pdf.pages)
        if count_words(text) < pages * 5:  # likely scanned
            ocr = try_ocr_pdf(fp)
            if ocr and count_words(ocr) > count_words(text):
                return truncate(ocr), pages, "ocr", True
        return truncate(text.strip()), pages, "native", False
    except ImportError:
        return None, 0, "need:pdfplumber", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_docx(fp):
    try:
        import docx
        doc = docx.Document(fp)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells)
        return truncate("\n".join(parts)), len(doc.paragraphs), "native", False
    except ImportError:
        return None, 0, "need:python-docx", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_doc_legacy(fp):
    """Legacy .doc via antiword or libreoffice fallback."""
    import subprocess
    for cmd in [
        ["antiword", str(fp)],
        ["libreoffice", "--headless", "--convert-to", "txt", "--outdir", "/tmp", str(fp)],
    ]:
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            if r.returncode == 0 and r.stdout.strip():
                return truncate(r.stdout), 0, f"cli:{cmd[0]}", False
            # libreoffice writes to file
            txt_path = Path("/tmp") / (Path(fp).stem + ".txt")
            if txt_path.exists():
                text = txt_path.read_text(errors="replace")
                txt_path.unlink()
                return truncate(text), 0, "libreoffice", False
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None, 0, "need:antiword|libreoffice", False


def extract_rtf(fp):
    try:
        from striprtf.striprtf import rtf_to_text
        raw = Path(fp).read_text(errors="replace")
        return truncate(rtf_to_text(raw)), 0, "native", False
    except ImportError:
        return None, 0, "need:striprtf", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_odt(fp):
    """OpenDocument text via zipfile XML parsing."""
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(fp) as z:
            with z.open("content.xml") as f:
                tree = ET.parse(f)
        ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
        parts = [elem.text or "" for elem in tree.iter()
                 if elem.tag.endswith("}p") or elem.tag.endswith("}h")]
        return truncate("\n".join(parts)), 0, "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def process_spreadsheet_rows(sheet_iterator):
    """Unified spreadsheet processor. sheet_iterator yields (name, row_iter) tuples."""
    parts = []
    for sheet_name, rows in sheet_iterator:
        parts.append(f"[Sheet: {sheet_name}]")
        row_count = 0
        for row in rows:
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
                row_count += 1
                if row_count >= DB_MAX_ROWS:
                    parts.append(f"[...truncated at {DB_MAX_ROWS} rows]")
                    break
    return truncate("\n".join(parts)), len(parts), "native", False


def extract_xlsx(fp):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(fp, read_only=True, data_only=True)
        def sheets():
            for name in wb.sheetnames:
                yield name, wb[name].iter_rows(values_only=True)
        result = process_spreadsheet_rows(sheets())
        wb.close()
        return result
    except ImportError:
        return None, 0, "need:openpyxl", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_xls(fp):
    try:
        import xlrd
        wb = xlrd.open_workbook(fp)
        def sheets():
            for name in wb.sheet_names():
                ws = wb.sheet_by_name(name)
                yield name, ([str(ws.cell_value(r, c)) for c in range(ws.ncols)]
                             for r in range(ws.nrows))
        return process_spreadsheet_rows(sheets())
    except ImportError:
        return None, 0, "need:xlrd", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_ods(fp):
    """OpenDocument Spreadsheet."""
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(fp) as z:
            with z.open("content.xml") as f:
                tree = ET.parse(f)
        tns = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
        txns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        def sheets():
            for table in tree.iter(f"{{{tns}}}table"):
                name = table.attrib.get(f"{{{tns}}}name", "?")
                def row_iter(tbl=table):
                    for row in tbl.iter(f"{{{tns}}}table-row"):
                        cells = []
                        for cell in row.iter(f"{{{tns}}}table-cell"):
                            texts = [p.text or "" for p in cell.iter(f"{{{txns}}}p")]
                            cells.append(" ".join(texts))
                        if any(c.strip() for c in cells):
                            yield cells
                yield name, row_iter()
        return process_spreadsheet_rows(sheets())
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_pptx(fp):
    try:
        from pptx import Presentation
        prs = Presentation(fp)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    parts.append(f"[Notes] {notes}")
        return truncate("\n".join(parts)), len(prs.slides), "native", False
    except ImportError:
        return None, 0, "need:python-pptx", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_csv_file(fp, delimiter=","):
    try:
        text = safe_read_text(fp)
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        parts = []
        for i, row in enumerate(reader):
            if i >= DB_MAX_ROWS:
                parts.append(f"[...truncated at {DB_MAX_ROWS} rows]")
                break
            parts.append(" | ".join(row))
        return truncate("\n".join(parts)), len(parts), "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_json_file(fp):
    try:
        text = safe_read_text(fp)
        data = json.loads(text)
        # For arrays, join items; for objects, format key-value
        if isinstance(data, list):
            parts = []
            for i, item in enumerate(data[:DB_MAX_ROWS]):
                if isinstance(item, dict):
                    parts.append(" | ".join(f"{k}: {v}" for k, v in item.items()))
                else:
                    parts.append(str(item))
            return truncate("\n".join(parts)), len(parts), "native", False
        else:
            return truncate(json.dumps(data, indent=2, default=str)), 0, "native", False
    except Exception:
        # Fall back to plain text
        return extract_text(fp)


def extract_jsonl(fp):
    try:
        text = safe_read_text(fp)
        parts = []
        for i, line in enumerate(text.strip().split("\n")):
            if i >= DB_MAX_ROWS:
                parts.append(f"[...truncated at {DB_MAX_ROWS} rows]")
                break
            try:
                obj = json.loads(line)
                if isinstance(obj, dict):
                    parts.append(" | ".join(f"{k}: {v}" for k, v in obj.items()))
                else:
                    parts.append(str(obj))
            except json.JSONDecodeError:
                parts.append(line)
        return truncate("\n".join(parts)), len(parts), "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_html(fp):
    try:
        raw = safe_read_text(fp)
        return truncate(strip_html(raw)), 0, "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_xml(fp):
    try:
        from xml.etree import ElementTree as ET
        tree = ET.parse(fp)
        parts = []
        for elem in tree.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return truncate("\n".join(parts)), 0, "native", False
    except Exception:
        return extract_text(fp)


def extract_sqlite(fp):
    """Extract schema + sample data from SQLite databases."""
    try:
        conn = sqlite3.connect(str(fp))
        conn.row_factory = sqlite3.Row
        parts = []

        tables = [r[0] for r in conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table' ORDER BY name"
        ).fetchall()]

        parts.append(f"SQLite Database: {len(tables)} tables")

        for table in tables:
            # Schema
            schema = conn.execute(
                f"SELECT sql FROM sqlite_master WHERE name=?", (table,)
            ).fetchone()
            if schema:
                parts.append(f"\n[Table: {table}]")
                parts.append(schema[0])

            # Row count
            try:
                count = conn.execute(f'SELECT COUNT(*) FROM "{table}"').fetchone()[0]
                parts.append(f"Rows: {count}")
            except Exception:
                count = 0

            # Sample data
            try:
                rows = conn.execute(
                    f'SELECT * FROM "{table}" LIMIT {min(DB_MAX_ROWS, 500)}'
                ).fetchall()
                if rows:
                    cols = [desc[0] for desc in conn.execute(f'SELECT * FROM "{table}" LIMIT 1').description]
                    parts.append(" | ".join(cols))
                    for row in rows:
                        parts.append(" | ".join(str(v) for v in row))
            except Exception:
                pass

        conn.close()
        return truncate("\n".join(parts)), len(tables), "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_access(fp):
    """MS Access via mdbtools (Linux) or pypyodbc."""
    import subprocess
    try:
        # Try mdbtools
        tables_raw = subprocess.run(
            ["mdb-tables", "-1", str(fp)], capture_output=True, text=True, timeout=30
        )
        if tables_raw.returncode == 0:
            tables = [t.strip() for t in tables_raw.stdout.strip().split("\n") if t.strip()]
            parts = [f"Access Database: {len(tables)} tables"]
            for table in tables[:50]:
                parts.append(f"\n[Table: {table}]")
                export = subprocess.run(
                    ["mdb-export", str(fp), table],
                    capture_output=True, text=True, timeout=30
                )
                if export.returncode == 0:
                    lines = export.stdout.strip().split("\n")
                    parts.extend(lines[:DB_MAX_ROWS])
            return truncate("\n".join(parts)), len(tables), "mdbtools", False
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None, 0, "need:mdbtools", False


def extract_eml(fp):
    try:
        import email
        from email import policy
        raw = Path(fp).read_bytes()
        msg = email.message_from_bytes(raw, policy=policy.default)
        parts = format_email_headers(msg.get('From', '?'), msg.get('To', '?'),
                                     msg.get('Date', '?'), msg.get('Subject', '?'))
        body = msg.get_body(preferencelist=("plain", "html"))
        if body:
            content = body.get_content()
            if body.get_content_type() == "text/html":
                content = strip_html(content)
            parts.append(content)
        for att in msg.iter_attachments():
            fname = att.get_filename()
            if fname:
                parts.append(f"[Attachment: {fname}]")
        return truncate("\n".join(parts)), 0, "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_msg(fp):
    try:
        import extract_msg
        msg = extract_msg.Message(str(fp))
        parts = format_email_headers(msg.sender or '?', msg.to or '?',
                                     msg.date or '?', msg.subject or '?')
        parts.append(msg.body or "")
        for att in (msg.attachments or []):
            parts.append(f"[Attachment: {att.longFilename or att.shortFilename}]")
        msg.close()
        return truncate("\n".join(parts)), 0, "native", False
    except ImportError:
        return None, 0, "need:extract-msg", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_mbox(fp):
    try:
        import mailbox
        mbox = mailbox.mbox(str(fp))
        parts = []
        count = 0
        for i, msg in enumerate(mbox):
            if i >= 500:
                parts.append(f"[...truncated at 500 messages]")
                break
            parts.append(f"\n--- Message {i+1} ---")
            parts.extend(format_email_headers(msg.get('From', '?'), msg.get('To', '?'),
                                              msg.get('Date', '?'), msg.get('Subject', '?')))
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            parts.append(payload.decode("utf-8", errors="replace")[:2000])
                            break
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode("utf-8", errors="replace")[:2000])
            count = i + 1
        mbox.close()
        return truncate("\n".join(parts)), count, "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_jupyter(fp):
    try:
        data = json.loads(Path(fp).read_text(encoding="utf-8"))
        parts = []
        for i, cell in enumerate(data.get("cells", [])):
            ctype = cell.get("cell_type", "?")
            source = "".join(cell.get("source", []))
            parts.append(f"[Cell {i+1} ({ctype})]")
            parts.append(source)
            # Include text outputs
            for output in cell.get("outputs", []):
                if "text" in output:
                    parts.append("".join(output["text"]))
        return truncate("\n".join(parts)), len(data.get("cells", [])), "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_image(fp):
    """OCR an image file."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(fp)
        text = pytesseract.image_to_string(img)
        return truncate(text.strip()) if text.strip() else None, 1, "ocr", True
    except ImportError:
        return None, 0, "need:pytesseract+Pillow", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


_ARCHIVE_TEXT_EXTS = {".txt", ".csv", ".json", ".md", ".xml", ".html", ".log", ".sql"}


def process_archive_members(members, read_fn, deep, archive_label):
    """Unified archive processing. members: list of (name, size). read_fn(name)->bytes."""
    parts = [f"{archive_label}: {len(members)} files"]
    if not deep:
        for name, size in members[:ARCHIVE_MAX_FILES]:
            parts.append(f"  {name} ({size:,} bytes)")
        return truncate("\n".join(parts)), len(members), "listing", False
    extracted = 0
    for name, size in members[:ARCHIVE_MAX_FILES]:
        ext = Path(name).suffix.lower()
        if ext in _ARCHIVE_TEXT_EXTS:
            try:
                data = read_fn(name)
                if data:
                    parts.append(f"\n[{name}]")
                    parts.append(data.decode("utf-8", errors="replace")[:5000])
                    extracted += 1
            except Exception:
                pass
    return truncate("\n".join(parts)), extracted if deep else len(members), "deep", False


def extract_archive_zip(fp, deep=False):
    try:
        with zipfile.ZipFile(fp) as z:
            members = [(n, z.getinfo(n).file_size) for n in z.namelist()]
            return process_archive_members(members, z.read, deep, "ZIP Archive")
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_archive_tar(fp, deep=False):
    try:
        with tarfile.open(fp) as t:
            raw_members = t.getmembers()
            members = [(m.name, m.size) for m in raw_members]
            def read_fn(name):
                m = t.getmember(name)
                if m.isfile():
                    f = t.extractfile(m)
                    return f.read() if f else None
                return None
            return process_archive_members(members, read_fn, deep, "TAR Archive")
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_ical(fp):
    try:
        text = safe_read_text(fp)
        # Simple extraction of SUMMARY, DTSTART, LOCATION, DESCRIPTION
        parts = []
        events = re.split(r"BEGIN:VEVENT", text)
        for event in events[1:]:
            summary = re.search(r"SUMMARY:(.*)", event)
            dtstart = re.search(r"DTSTART[^:]*:(.*)", event)
            location = re.search(r"LOCATION:(.*)", event)
            desc = re.search(r"DESCRIPTION:(.*)", event)
            parts.append(f"Event: {summary.group(1).strip() if summary else '?'}")
            if dtstart:
                parts.append(f"  Date: {dtstart.group(1).strip()}")
            if location:
                parts.append(f"  Location: {location.group(1).strip()}")
            if desc:
                parts.append(f"  {desc.group(1).strip()[:200]}")
        return truncate("\n".join(parts)), len(events) - 1, "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_vcard(fp):
    try:
        text = safe_read_text(fp)
        parts = []
        cards = re.split(r"BEGIN:VCARD", text)
        for card in cards[1:]:
            fn = re.search(r"FN:(.*)", card)
            email = re.search(r"EMAIL[^:]*:(.*)", card)
            tel = re.search(r"TEL[^:]*:(.*)", card)
            org = re.search(r"ORG:(.*)", card)
            parts.append(f"Contact: {fn.group(1).strip() if fn else '?'}")
            if org:
                parts.append(f"  Org: {org.group(1).strip()}")
            if email:
                parts.append(f"  Email: {email.group(1).strip()}")
            if tel:
                parts.append(f"  Phone: {tel.group(1).strip()}")
        return truncate("\n".join(parts)), len(cards) - 1, "native", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def extract_text(fp):
    try:
        text = safe_read_text(fp)
        lines = text.count("\n") + 1
        return truncate(text), lines, "text_read", False
    except Exception as e:
        return None, 0, f"err:{e}"[:80], False


def try_ocr_pdf(filepath):
    """OCR a scanned PDF using tesseract + pdftoppm."""
    try:
        import subprocess, tempfile
        subprocess.run(["tesseract", "--version"], capture_output=True, check=True)
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["pdftoppm", "-png", "-r", "300", filepath, os.path.join(tmpdir, "pg")],
                capture_output=True, timeout=120
            )
            texts = []
            for img in sorted(Path(tmpdir).glob("*.png")):
                r = subprocess.run(
                    ["tesseract", str(img), "stdout"], capture_output=True, text=True, timeout=60
                )
                if r.returncode == 0:
                    texts.append(r.stdout)
            return "\n".join(texts) if texts else None
    except Exception:
        return None


# ─── Extraction Router ────────────────────────────────────────────────────────

EXTRACTORS = {
    "pdf": extract_pdf, "docx": extract_docx, "doc_legacy": extract_doc_legacy,
    "rtf": extract_rtf, "odt": extract_odt,
    "xlsx": extract_xlsx, "xls": extract_xls, "xlsb": extract_xlsx, "ods": extract_ods,
    "pptx": extract_pptx,
    "csv": lambda fp: extract_csv_file(fp, ","),
    "tsv": lambda fp: extract_csv_file(fp, "\t"),
    "json_file": extract_json_file, "jsonl": extract_jsonl,
    "html": extract_html, "xml": extract_xml,
    "sqlite": extract_sqlite, "access": extract_access,
    "eml": extract_eml, "msg": extract_msg, "mbox": extract_mbox,
    "jupyter": extract_jupyter,
    "image": extract_image,
    "archive_zip": extract_archive_zip, "archive_tar": extract_archive_tar,
    "archive_gz": extract_archive_tar,
    "ical": extract_ical, "vcard": extract_vcard,
    "text": extract_text,
}

def extract_file(filepath, deep=False):
    ext = Path(filepath).suffix.lower()
    doc_type = SUPPORTED_TYPES.get(ext)
    if not doc_type:
        return None, 0, "unsupported", False
    extractor = EXTRACTORS.get(doc_type, extract_text)
    if doc_type.startswith("archive_") and deep:
        return extractor(filepath, deep=True)
    return extractor(filepath)


# ─── Entity Extraction ────────────────────────────────────────────────────────

# Pre-compiled entity patterns (Phase 4B: compiled once at module load)
ENTITY_PATTERNS = [
    ("date",        re.compile(r'\b(\d{1,2}/\d{1,2}/\d{2,4})\b', re.IGNORECASE), None),
    ("date",        re.compile(r'\b(\d{1,2}-\d{1,2}-\d{2,4})\b', re.IGNORECASE), None),
    ("date",        re.compile(r'\b((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{1,2},? \d{4})\b', re.IGNORECASE), None),
    ("date",        re.compile(r'\b(\d{4}-\d{2}-\d{2})\b'), None),
    ("amount",      re.compile(r'\$[\d,]+(?:\.\d{2})?'), None),
    ("email",       re.compile(r'\b[\w.+-]+@[\w-]+\.[\w.-]+\b'), None),
    ("phone",       re.compile(r'\b(?:\+?1[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b'), None),
    ("case_number", re.compile(r'\b(?:Case|Docket|No\.|Index|File)[\s#:]*[\d-]+[A-Za-z]*/?\d*\b', re.IGNORECASE), None),
    ("ssn_detected",re.compile(r'\b\d{3}-\d{2}-\d{4}\b'), lambda m: "***-**-" + m.group()[-4:]),
    ("url",         re.compile(r'https?://[^\s<>"]+'), lambda m: m.group()[:200]),
    ("address",     re.compile(r'\b\d+\s+[\w\s]+(?:Street|St|Avenue|Ave|Road|Rd|Drive|Dr|Boulevard|Blvd|Lane|Ln|Way|Court|Ct|Place|Pl)\b', re.IGNORECASE),
                    lambda m: m.group().strip() if len(m.group()) > 8 else None),
]


def extract_entities(text):
    entities = []
    if not text:
        return entities
    for etype, pattern, transform in ENTITY_PATTERNS:
        for m in pattern.finditer(text):
            if transform:
                value = transform(m)
                if value is None:
                    continue
            else:
                value = m.group(1) if pattern.groups else m.group()
            entities.append((etype, value, get_ctx(text, m.start())))
    return entities


def get_ctx(text, pos, window=80):
    start = max(0, pos - window)
    end = min(len(text), pos + window)
    return text[start:end].replace("\n", " ").strip()


# ─── Commands ─────────────────────────────────────────────────────────────────

def walk_supported_files(repo_path):
    """Generator-based file walking (Phase 5A: O(1) memory for file list)."""
    skip_dirs = {".git", "node_modules", "__pycache__", ".venv", "venv", "dist", "build",
                 ".next", "vendor", "env", ".tox", ".mypy_cache", ".pytest_cache",
                 ".eggs", "egg-info", "__MACOSX", ".DS_Store"}
    for root, dirs, fnames in os.walk(repo_path):
        dirs[:] = [d for d in dirs if d not in skip_dirs and not d.startswith(".")]
        for fname in fnames:
            full = Path(root) / fname
            try:
                if full.stat().st_size > MAX_FILE_SIZE:
                    continue
            except OSError:
                continue
            ext = full.suffix.lower()
            if ext in SUPPORTED_TYPES:
                yield full, str(full.relative_to(repo_path))


def cmd_scan(project_name, deep=False):
    conn = get_db()
    conn.executescript(DOC_SCHEMA)

    repo = conn.execute("SELECT id, path, name FROM repos WHERE name = ?", (project_name,)).fetchone()
    if not repo:
        print(f"[ERROR] Project '{project_name}' not found in registry.")
        conn.close()
        return

    repo_id, repo_path_str, repo_name = repo["id"], repo["path"], repo["name"]
    repo_path = Path(repo_path_str)

    if not repo_path.exists():
        print(f"[ERROR] Path does not exist: {repo_path}")
        conn.close()
        return

    # Count files first for progress (lightweight pass)
    files = list(walk_supported_files(repo_path))
    total = len(files)

    print(f"\n  Scanning {total} files in '{project_name}'...")
    print(f"  Path: {repo_path}")
    print(f"  Deep mode: {'ON' if deep else 'OFF (use --deep for archives/DBs)'}")
    print(f"  {'=' * 55}")

    stats = {"processed": 0, "skipped": 0, "ocrd": 0, "errors": 0,
             "words": 0, "entities": 0, "categories": {}}

    for full_path, rel_path in files:
        fname = full_path.name
        fsize = full_path.stat().st_size
        fhash = file_hash(str(full_path))
        ext = full_path.suffix.lower()
        doc_type = SUPPORTED_TYPES.get(ext, "unknown")
        category = CATEGORY_MAP.get(doc_type, "other")

        existing = conn.execute(
            "SELECT file_hash FROM documents WHERE repo_id = ? AND file_path = ?",
            (repo_id, rel_path)
        ).fetchone()

        if existing and existing["file_hash"] == fhash:
            stats["skipped"] += 1
            continue

        text, page_count, method, is_ocr = extract_file(str(full_path), deep=deep)

        wc = count_words(text) if text else 0
        cc = len(text) if text else 0
        stats["words"] += wc
        if is_ocr:
            stats["ocrd"] += 1
        if text is None and method and method.startswith("err"):
            stats["errors"] += 1

        stats["categories"][category] = stats["categories"].get(category, 0) + 1

        mod_date = datetime.fromtimestamp(full_path.stat().st_mtime).isoformat()

        # Use cursor.lastrowid instead of separate SELECT (Phase 5E)
        cursor = conn.execute("""
            INSERT INTO documents
                (repo_id, file_path, file_name, file_type, file_category, file_hash, file_size,
                 page_count, word_count, char_count, extracted_text, extraction_method,
                 is_searchable, is_ocrd, modified_date, scanned_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(repo_id, file_path) DO UPDATE SET
                file_hash=excluded.file_hash, file_size=excluded.file_size,
                file_category=excluded.file_category,
                page_count=excluded.page_count, word_count=excluded.word_count,
                char_count=excluded.char_count, extracted_text=excluded.extracted_text,
                extraction_method=excluded.extraction_method,
                is_searchable=excluded.is_searchable, is_ocrd=excluded.is_ocrd,
                modified_date=excluded.modified_date, scanned_at=excluded.scanned_at
        """, (repo_id, rel_path, fname, ext, category, fhash, fsize,
              page_count, wc, cc, text, method,
              1 if text and wc > 0 else 0, 1 if is_ocr else 0,
              mod_date, now_iso()))

        if text:
            # Use lastrowid for inserts, fall back to SELECT for upsert updates
            doc_id = cursor.lastrowid
            if not doc_id:
                doc_id = conn.execute(
                    "SELECT id FROM documents WHERE repo_id = ? AND file_path = ?",
                    (repo_id, rel_path)
                ).fetchone()["id"]

            conn.execute("DELETE FROM doc_search WHERE doc_id = ?", (doc_id,))
            conn.execute(
                """INSERT INTO doc_search(file_path, file_name, extracted_text, repo_name, repo_id, doc_id)
                   VALUES (?, ?, ?, ?, ?, ?)""",
                (rel_path, fname, text, repo_name, repo_id, doc_id)
            )

            # Batch entity insertion (Phase 5D)
            conn.execute("DELETE FROM doc_entities WHERE doc_id = ?", (doc_id,))
            entity_list = extract_entities(text)
            if entity_list:
                conn.executemany(
                    "INSERT INTO doc_entities (doc_id, repo_id, entity_type, entity_value, context) VALUES (?,?,?,?,?)",
                    [(doc_id, repo_id, etype, evalue, ectx) for etype, evalue, ectx in entity_list]
                )
                stats["entities"] += len(entity_list)

        stats["processed"] += 1
        if total > 0:
            pct = ((stats["processed"] + stats["skipped"]) / total) * 100
            tag = "OCR" if is_ocr else category[:6]
            print(f"  [{pct:5.1f}%] {tag:<8} {rel_path[:65]}")

        # Periodic commits every 100 files (Phase 5F)
        if (stats["processed"] + stats["skipped"]) % 100 == 0:
            conn.commit()

    conn.commit()

    print(f"""
  {'=' * 55}
  SCAN COMPLETE: {project_name}
  {'=' * 55}
  Processed:    {stats['processed']}
  Skipped:      {stats['skipped']} (unchanged)
  OCR'd:        {stats['ocrd']}
  Errors:       {stats['errors']}
  Total Words:  {stats['words']:,}
  Entities:     {stats['entities']:,}

  By Category:""")
    for cat, cnt in sorted(stats["categories"].items(), key=lambda x: -x[1]):
        print(f"    {cat:<15} {cnt:>4} files")
    print(f"  {'=' * 55}\n")

    conn.close()


def cmd_search(query, project_name=None):
    with db_session(DOC_SCHEMA) as conn:
        if project_name:
            repo = conn.execute("SELECT id FROM repos WHERE name = ?", (project_name,)).fetchone()
            if not repo:
                print(f"[ERROR] Project '{project_name}' not found.")
                return
            results = conn.execute("""
                SELECT file_path, file_name, repo_name, repo_id, doc_id,
                       snippet(doc_search, 2, '>>>', '<<<', '...', 40) as snippet
                FROM doc_search
                WHERE doc_search MATCH ? AND repo_id = CAST(? AS TEXT)
                ORDER BY rank LIMIT 20
            """, (query, str(repo["id"]))).fetchall()
        else:
            results = conn.execute("""
                SELECT file_path, file_name, repo_name, repo_id, doc_id,
                       snippet(doc_search, 2, '>>>', '<<<', '...', 40) as snippet
                FROM doc_search
                WHERE doc_search MATCH ?
                ORDER BY rank LIMIT 20
            """, (query,)).fetchall()

        if not results:
            print(f"  No results for '{query}'")
            return

        print(f"\n  SEARCH: '{query}' ({len(results)} matches)")
        print(f"  {'=' * 60}")
        current_repo = None
        for r in results:
            if r["repo_name"] != current_repo:
                current_repo = r["repo_name"]
                print(f"\n  [{current_repo}]")
            snippet = r["snippet"].replace("\n", " ").strip()
            print(f"    {r['file_path']}")
            print(f"      ...{snippet}...")
            print()


def cmd_info(project_name):
    with db_session(DOC_SCHEMA) as conn:
        repo = conn.execute("SELECT id FROM repos WHERE name = ?", (project_name,)).fetchone()
        if not repo:
            print(f"[ERROR] Project '{project_name}' not found.")
            return

        rid = repo["id"]
        s = conn.execute("""
            SELECT COUNT(*) as total, SUM(word_count) as words, SUM(file_size) as size,
                   SUM(page_count) as pages, SUM(is_ocrd) as ocrd,
                   SUM(CASE WHEN is_searchable=1 THEN 1 ELSE 0 END) as searchable,
                   SUM(CASE WHEN is_searchable=0 THEN 1 ELSE 0 END) as unsearchable
            FROM documents WHERE repo_id = ?
        """, (rid,)).fetchone()

        if not s or s["total"] == 0:
            print(f"  No documents indexed for '{project_name}'. Run: doc-intel.py scan {project_name}")
            return

        cats = conn.execute("""
            SELECT file_category, COUNT(*) as cnt, SUM(word_count) as words, SUM(file_size) as size
            FROM documents WHERE repo_id = ? GROUP BY file_category ORDER BY cnt DESC
        """, (rid,)).fetchall()

        types = conn.execute("""
            SELECT file_type, COUNT(*) as cnt
            FROM documents WHERE repo_id = ? GROUP BY file_type ORDER BY cnt DESC
        """, (rid,)).fetchall()

        entities = conn.execute("""
            SELECT entity_type, COUNT(*) as cnt
            FROM doc_entities WHERE repo_id = ? GROUP BY entity_type ORDER BY cnt DESC
        """, (rid,)).fetchall()

        methods = conn.execute("""
            SELECT extraction_method, COUNT(*) as cnt
            FROM documents WHERE repo_id = ? GROUP BY extraction_method ORDER BY cnt DESC
        """, (rid,)).fetchall()

        largest = conn.execute("""
            SELECT file_path, word_count, file_category, file_size
            FROM documents WHERE repo_id = ? AND word_count > 0
            ORDER BY word_count DESC LIMIT 10
        """, (rid,)).fetchall()

        mb = (s["size"] or 0) / (1024 * 1024)

        print(f"""
  +{'=' * 58}+
  : DOCUMENT INTELLIGENCE: {project_name:<33}:
  +{'=' * 58}+

  Total Files:       {s['total']}
  Searchable:        {s['searchable']}
  Not Searchable:    {s['unsearchable']}
  OCR'd:             {s['ocrd'] or 0}
  Total Pages:       {s['pages'] or 0:,}
  Total Words:       {s['words'] or 0:,}
  Total Size:        {mb:.1f} MB""")

        if cats:
            print(f"\n  BY CATEGORY\n  {'=' * 40}")
            for c in cats:
                sz = (c["size"] or 0) / (1024 * 1024)
                print(f"    {c['file_category'] or '?':<15} {c['cnt']:>4} files  {c['words'] or 0:>8,} words  {sz:.1f} MB")

        if types:
            print(f"\n  BY FILE TYPE\n  {'=' * 40}")
            for t in types:
                print(f"    {t['file_type'] or '?':<10} {t['cnt']:>4}")

        if entities:
            print(f"\n  ENTITIES FOUND\n  {'=' * 40}")
            for e in entities:
                print(f"    {e['entity_type']:<15} {e['cnt']:>5}")

        if methods:
            print(f"\n  EXTRACTION METHODS\n  {'=' * 40}")
            for m in methods:
                print(f"    {m['extraction_method'] or '?':<20} {m['cnt']:>4} files")

        if largest:
            print(f"\n  LARGEST (by words)\n  {'=' * 40}")
            for d in largest:
                print(f"    {d['word_count']:>8,} words  [{d['file_category']}]  {d['file_path']}")

        print(f"\n  +{'=' * 58}+")


def cmd_entities(project_name, entity_type=None):
    with db_session(DOC_SCHEMA) as conn:
        repo = conn.execute("SELECT id FROM repos WHERE name = ?", (project_name,)).fetchone()
        if not repo:
            print(f"[ERROR] Project '{project_name}' not found.")
            return

        if entity_type:
            rows = conn.execute("""
                SELECT e.entity_type, e.entity_value, d.file_path
                FROM doc_entities e JOIN documents d ON d.id = e.doc_id
                WHERE e.repo_id = ? AND e.entity_type = ? ORDER BY e.entity_value
            """, (repo["id"], entity_type)).fetchall()
        else:
            rows = conn.execute("""
                SELECT e.entity_type, e.entity_value, d.file_path
                FROM doc_entities e JOIN documents d ON d.id = e.doc_id
                WHERE e.repo_id = ? ORDER BY e.entity_type, e.entity_value
            """, (repo["id"],)).fetchall()

        if not rows:
            print(f"  No entities found for '{project_name}'.")
            return

        print(f"\n  ENTITIES: {project_name} ({len(rows)} total)")
        print(f"  {'=' * 50}")
        current = None
        for r in rows:
            if r["entity_type"] != current:
                current = r["entity_type"]
                print(f"\n  [{current.upper()}]")
            print(f"    {r['entity_value']:<35} in {r['file_path']}")


def cmd_summary():
    """Cross-project document intelligence overview."""
    with db_session(DOC_SCHEMA) as conn:
        totals = conn.execute("""
            SELECT COUNT(*) as docs, SUM(word_count) as words, SUM(file_size) as size,
                   COUNT(DISTINCT repo_id) as projects, SUM(is_ocrd) as ocrd
            FROM documents
        """).fetchone()

        if not totals or totals["docs"] == 0:
            print("  No documents indexed yet. Run: doc-intel.py scan <project-name>")
            return

        by_project = conn.execute("""
            SELECT r.name, COUNT(*) as docs, SUM(d.word_count) as words,
                   SUM(d.file_size) as size,
                   SUM(CASE WHEN d.is_searchable=1 THEN 1 ELSE 0 END) as searchable
            FROM documents d JOIN repos r ON r.id = d.repo_id
            GROUP BY r.name ORDER BY docs DESC
        """).fetchall()

        by_cat = conn.execute("""
            SELECT file_category, COUNT(*) as cnt, SUM(word_count) as words
            FROM documents GROUP BY file_category ORDER BY cnt DESC
        """).fetchall()

        entity_totals = conn.execute("""
            SELECT entity_type, COUNT(*) as cnt
            FROM doc_entities GROUP BY entity_type ORDER BY cnt DESC
        """).fetchall()

        mb = (totals["size"] or 0) / (1024 * 1024)
        print(f"""
  +{'=' * 58}+
  :            DOCUMENT INTELLIGENCE OVERVIEW                   :
  +{'=' * 58}+

  Total Documents:   {totals['docs']:,}
  Total Words:       {totals['words'] or 0:,}
  Total Size:        {mb:.1f} MB
  Projects Indexed:  {totals['projects']}
  OCR'd Files:       {totals['ocrd'] or 0}""")

        if by_project:
            print(f"\n  BY PROJECT\n  {'=' * 50}")
            for p in by_project:
                sz = (p["size"] or 0) / (1024 * 1024)
                print(f"    {p['name']:<25} {p['docs']:>4} files  {p['words'] or 0:>8,} words  {sz:.1f} MB")

        if by_cat:
            print(f"\n  BY CATEGORY\n  {'=' * 50}")
            for c in by_cat:
                print(f"    {c['file_category'] or '?':<15} {c['cnt']:>5} files  {c['words'] or 0:>8,} words")

        if entity_totals:
            print(f"\n  ALL ENTITIES\n  {'=' * 50}")
            for e in entity_totals:
                print(f"    {e['entity_type']:<15} {e['cnt']:>6}")

        print(f"\n  +{'=' * 58}+")


def cmd_deps_check():
    print(f"\n  EXTRACTION DEPENDENCY CHECK")
    print(f"  {'=' * 50}")

    libs = {
        "pdfplumber":   ("pip install pdfplumber",        "PDF text extraction"),
        "docx":         ("pip install python-docx",       "Word .docx extraction"),
        "openpyxl":     ("pip install openpyxl",          "Excel .xlsx extraction"),
        "pptx":         ("pip install python-pptx",       "PowerPoint .pptx extraction"),
        "xlrd":         ("pip install xlrd",              "Legacy Excel .xls extraction"),
        "striprtf":     ("pip install striprtf",          "RTF extraction"),
        "extract_msg":  ("pip install extract-msg",       "Outlook .msg email extraction"),
        "pytesseract":  ("pip install pytesseract Pillow","Image OCR (+ Pillow for images)"),
        "chardet":      ("pip install chardet",           "Encoding detection for text files"),
    }

    installed = 0
    for lib, (install_cmd, desc) in libs.items():
        try:
            __import__(lib)
            print(f"  [OK]    {lib:<15} {desc}")
            installed += 1
        except ImportError:
            print(f"  [MISS]  {lib:<15} {desc}")
            print(f"          {install_cmd}")

    print(f"\n  System tools:")
    import subprocess
    for tool, desc, install in [
        ("tesseract", "OCR engine for scanned docs/images",
         "choco install tesseract / brew install tesseract / apt install tesseract-ocr"),
        ("pdftoppm", "PDF to image for OCR pipeline",
         "choco install poppler / brew install poppler / apt install poppler-utils"),
        ("antiword", "Legacy .doc extraction",
         "choco install antiword / apt install antiword"),
        ("mdb-tables", "MS Access .mdb extraction",
         "apt install mdbtools (Linux only)"),
        ("libreoffice", "Fallback for .doc .odt .ods .odp",
         "choco install libreoffice / brew install --cask libreoffice"),
    ]:
        try:
            subprocess.run([tool, "--version"], capture_output=True, timeout=5)
            print(f"  [OK]    {tool:<15} {desc}")
            installed += 1
        except (FileNotFoundError, subprocess.TimeoutExpired):
            print(f"  [MISS]  {tool:<15} {desc}")
            print(f"          {install}")

    total = len(libs) + 5
    print(f"\n  {installed}/{total} extractors available.")
    print(f"  Core 4 (pdfplumber, python-docx, openpyxl, python-pptx) cover 90% of use cases.\n")


def cmd_rebuild():
    with db_session(DOC_SCHEMA) as conn:
        conn.execute("DELETE FROM doc_search")
        conn.execute("""
            INSERT INTO doc_search(file_path, file_name, extracted_text, repo_name, repo_id, doc_id)
            SELECT d.file_path, d.file_name, d.extracted_text, r.name, d.repo_id, d.id
            FROM documents d JOIN repos r ON r.id = d.repo_id
            WHERE d.is_searchable = 1
        """)
        count = conn.execute("SELECT COUNT(*) as c FROM doc_search").fetchone()["c"]
        print(f"[OK] Search index rebuilt. {count} searchable documents.")


# ─── CLI Router ───────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        return

    cmd = sys.argv[1].lower()

    if cmd == "scan":
        if len(sys.argv) < 3:
            print("Usage: doc-intel.py scan <project-name> [--deep]")
            return
        deep = "--deep" in sys.argv
        cmd_scan(sys.argv[2], deep=deep)

    elif cmd == "search":
        if len(sys.argv) < 3:
            print("Usage: doc-intel.py search <query> [--project <n>]")
            return
        query = sys.argv[2]
        project = sys.argv[sys.argv.index("--project") + 1] if "--project" in sys.argv else None
        cmd_search(query, project)

    elif cmd == "info":
        if len(sys.argv) < 3:
            return print("Usage: doc-intel.py info <project-name>")
        cmd_info(sys.argv[2])

    elif cmd == "entities":
        if len(sys.argv) < 3:
            return print("Usage: doc-intel.py entities <project-name> [type]")
        cmd_entities(sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else None)

    elif cmd == "summary":
        cmd_summary()

    elif cmd == "deps-check":
        cmd_deps_check()

    elif cmd == "rebuild":
        cmd_rebuild()

    else:
        print(f"Unknown command: {cmd}")
        print(__doc__)


if __name__ == "__main__":
    main()
